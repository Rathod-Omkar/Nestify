import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    # Set widescreen 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette Definitions (Nestify Brand Theme)
    DARK_BLUE = RGBColor(16, 28, 48)     # Primary Brand Color (Deep Navy)
    TEAL = RGBColor(0, 180, 162)         # Accent Color (Teal)
    LIGHT_GRAY = RGBColor(245, 246, 248) # Slide Background Color (Light Gray)
    DARK_GRAY = RGBColor(50, 50, 50)     # Body Text Color (Slate)
    WHITE = RGBColor(255, 255, 255)      # White
    LIGHT_BLUE = RGBColor(235, 245, 245) # Card Fill Color (Very Light Teal/Blue)
    ACCENT_RED = RGBColor(194, 24, 91)   # Attention/Alert Color (Fuchsia)

    # Reusable slide helpers
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text):
        # Top banner shape
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = DARK_BLUE
        header_shape.line.fill.background() # No border
        
        # Header text
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.15), Inches(11.833), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.name = "Calibri"
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Underline accent band
        accent_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(13.333), Inches(0.08))
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = TEAL
        accent_shape.line.fill.background()

    # SLIDE 1: Title Slide (Dark Theme)
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide1, DARK_BLUE)

    # Decorative Teal band on the left
    band = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = TEAL
    band.line.fill.background()

    # Title text box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.333), Inches(2.4))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "NESTIFY"
    p1.font.name = "Calibri"
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_after = Pt(10)

    p2 = tf1.add_paragraph()
    p2.text = "A Community-Driven Shared Accommodation Booking Platform"
    p2.font.name = "Calibri"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = TEAL
    p2.space_after = Pt(5)

    p3 = tf1.add_paragraph()
    p3.text = "MERN Stack Web Application solving Travel Costs, Solo Female Safety, and Eco-Friendly Feature Disclosures in Lodging"
    p3.font.name = "Calibri"
    p3.font.size = Pt(16)
    p3.font.italic = True
    p3.font.color.rgb = LIGHT_GRAY

    # Student metadata details
    student_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(5.5), Inches(2.0))
    tfs = student_box.text_frame
    tfs.word_wrap = True
    ps_head = tfs.paragraphs[0]
    ps_head.text = "Submitted By:"
    ps_head.font.name = "Calibri"
    ps_head.font.size = Pt(14)
    ps_head.font.bold = True
    ps_head.font.color.rgb = TEAL
    ps_head.space_after = Pt(6)

    details = [
        "Name: [Your Name]",
        "Class & Div: [Class & Div]",
        "Roll No: [Roll No]"
    ]
    for det in details:
        p = tfs.add_paragraph()
        p.text = det
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY

    # Guide metadata details
    guide_box = slide1.shapes.add_textbox(Inches(7.0), Inches(4.8), Inches(5.5), Inches(2.0))
    tfg = guide_box.text_frame
    tfg.word_wrap = True
    pg_head = tfg.paragraphs[0]
    pg_head.text = "Under the Guidance of:"
    pg_head.font.name = "Calibri"
    pg_head.font.size = Pt(14)
    pg_head.font.bold = True
    pg_head.font.color.rgb = TEAL
    pg_head.space_after = Pt(6)

    guide_details = [
        "Guide: [Guide Name]",
        "Department of Computer Engineering",
        "Institution Name"
    ]
    for gdet in guide_details:
        p = tfg.add_paragraph()
        p.text = gdet
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY

    # SLIDE 2: Project Completion Letter
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide2, LIGHT_GRAY)
    add_header(slide2, "Project Completion Letter")

    # Add Letter Body Card
    card2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.3))
    card2.fill.solid()
    card2.fill.fore_color.rgb = WHITE
    card2.line.color.rgb = TEAL
    card2.line.width = Pt(1.5)

    letter_box = slide2.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(11.333), Inches(4.9))
    tfl = letter_box.text_frame
    tfl.word_wrap = True
    
    pl1 = tfl.paragraphs[0]
    pl1.text = "CERTIFICATE OF COMPLETION"
    pl1.alignment = PP_ALIGN.CENTER
    pl1.font.name = "Calibri"
    pl1.font.size = Pt(22)
    pl1.font.bold = True
    pl1.font.color.rgb = DARK_BLUE
    pl1.space_after = Pt(16)

    pl2 = tfl.add_paragraph()
    pl2.text = (
        "This is to certify that the project report entitled \"NESTIFY: A Community-Driven Shared "
        "Accommodation Booking Platform\" is a bonafide work carried out by [Your Name], student of "
        "[Class & Div], bearing Roll No: [Roll No], towards the partial fulfillment of the requirements for "
        "the academic curriculum.\n\n"
        "The project has been successfully completed in accordance with the department guidelines and is approved "
        "for submission. All design documentation, implementation modules, and testing schedules have been executed "
        "to a satisfactory academic standard under my guidance and supervision."
    )
    pl2.font.name = "Calibri"
    pl2.font.size = Pt(15)
    pl2.font.color.rgb = DARK_GRAY
    pl2.space_after = Pt(40)

    pl_signs = tfl.add_paragraph()
    pl_signs.text = (
        "_________________________\t\t\t\t_________________________\n"
        "Internal Guide / Supervisor\t\t\t\tHead of Department\n\n\n"
        "_________________________\t\t\t\t_________________________\n"
        "Internal Examiner\t\t\t\tExternal Examiner"
    )
    pl_signs.font.name = "Calibri"
    pl_signs.font.size = Pt(14)
    pl_signs.font.bold = True
    pl_signs.font.color.rgb = DARK_BLUE

    # SLIDE 3: Introduction (Project Overview + Problem Gaps)
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide3, LIGHT_GRAY)
    add_header(slide3, "Introduction")

    # Column 1: What is Nestify?
    col1_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), Inches(5.7), Inches(5.3))
    col1_card.fill.solid()
    col1_card.fill.fore_color.rgb = WHITE
    col1_card.line.fill.background()
    
    col1_box = slide3.shapes.add_textbox(Inches(0.95), Inches(1.7), Inches(5.3), Inches(4.9))
    tfc1 = col1_box.text_frame
    tfc1.word_wrap = True
    
    pc1_h = tfc1.paragraphs[0]
    pc1_h.text = "What is Nestify?"
    pc1_h.font.name = "Calibri"
    pc1_h.font.size = Pt(18)
    pc1_h.font.bold = True
    pc1_h.font.color.rgb = DARK_BLUE
    pc1_h.space_after = Pt(10)

    # Project overview paragraph
    intro_para = tfc1.add_paragraph()
    intro_para.text = (
        "Nestify is a community-driven shared accommodation booking platform "
        "built on the MERN stack (MongoDB, Express.js, EJS, Node.js). It connects "
        "budget travelers with local homeowners who have spare rooms to share, "
        "creating a trust-based economy for affordable lodging."
    )
    intro_para.font.name = "Calibri"
    intro_para.font.size = Pt(12.5)
    intro_para.font.color.rgb = DARK_GRAY
    intro_para.space_after = Pt(12)

    # Tech stack heading
    tech_h = tfc1.add_paragraph()
    tech_h.text = "Technology Stack"
    tech_h.font.name = "Calibri"
    tech_h.font.size = Pt(15)
    tech_h.font.bold = True
    tech_h.font.color.rgb = TEAL
    tech_h.space_after = Pt(6)

    tech_items = [
        "Backend: Node.js + Express.js (MVC Architecture)",
        "Database: MongoDB with Mongoose ODM",
        "Frontend: EJS Templates + Bootstrap 5 + Flatpickr",
        "Auth: Passport.js with session-based login",
        "Uploads: Multer for images & document verification"
    ]
    for item in tech_items:
        p = tfc1.add_paragraph()
        p.text = f"• {item}"
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

    # Key features heading
    feat_h = tfc1.add_paragraph()
    feat_h.text = "Key Features"
    feat_h.font.name = "Calibri"
    feat_h.font.size = Pt(15)
    feat_h.font.bold = True
    feat_h.font.color.rgb = TEAL
    feat_h.space_after = Pt(6)

    features = [
        "Admin-vetted listings with document verification",
        "Date collision prevention for bookings",
        "Verified women's safety ratings for female travelers",
        "Eco-friendly feature badges on property listings",
        "Community services marketplace (guides, drivers, helpers)"
    ]
    for feat in features:
        p = tfc1.add_paragraph()
        p.text = f"✔ {feat}"
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

    # Column 2: Problems Nestify Solves
    col2_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.883), Inches(1.5), Inches(5.7), Inches(5.3))
    col2_card.fill.solid()
    col2_card.fill.fore_color.rgb = WHITE
    col2_card.line.fill.background()
    
    col2_box = slide3.shapes.add_textbox(Inches(7.083), Inches(1.7), Inches(5.3), Inches(4.9))
    tfc2 = col2_box.text_frame
    tfc2.word_wrap = True

    pc2_h = tfc2.paragraphs[0]
    pc2_h.text = "Problems Nestify Addresses"
    pc2_h.font.name = "Calibri"
    pc2_h.font.size = Pt(18)
    pc2_h.font.bold = True
    pc2_h.font.color.rgb = DARK_BLUE
    pc2_h.space_after = Pt(10)

    problems2 = [
        ("High Travel Costs:", " Mainstream aggregators prioritize premium hotels and villas. Budget travelers, students, and backpackers struggle to find affordable shared spaces with a home-like atmosphere."),
        ("Underutilized Rooms:", " Homeowners with spare rooms have no trusted, simple platform to list and monetize these spaces while welcoming cultural exchanges with travelers."),
        ("Female Safety Gaps:", " Generic reviews on existing platforms do not isolate safety parameters. Solo female travelers need verified reviews from women who actually stayed at the property."),
        ("Listing Fraud:", " Many platforms allow hosts to list properties without identity verification, leading to fake or misleading listings that erode traveler trust."),
        ("No Eco Transparency:", " Eco-conscious travelers have no standard way to check which green amenities (solar energy, recycling, rainwater harvesting) a property actually has.")
    ]
    for title, desc in problems2:
        p = tfc2.add_paragraph()
        p.text = f"• {title}{desc}"
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(7)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = DARK_BLUE

    # SLIDE 4: Objectives of the Project (Detailed)
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide4, LIGHT_GRAY)
    add_header(slide4, "Project Objectives & Technical Goals")

    obj_card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.3))
    obj_card.fill.solid()
    obj_card.fill.fore_color.rgb = WHITE
    obj_card.line.fill.background()

    obj_box = slide4.shapes.add_textbox(Inches(0.95), Inches(1.6), Inches(11.333), Inches(5.1))
    tfo = obj_box.text_frame
    tfo.word_wrap = True

    po_h = tfo.paragraphs[0]
    po_h.text = "Concrete Project Targets & Expected Solutions"
    po_h.font.name = "Calibri"
    po_h.font.size = Pt(18)
    po_h.font.bold = True
    po_h.font.color.rgb = DARK_BLUE
    po_h.space_after = Pt(10)

    objectives = [
        ("Deployment of MVC Architecture:", " Build a web application using the Model-View-Controller design pattern (Mongoose Models, EJS-mate templates as Views, and Express.js Routers as Controllers) for a clean separation of concern."),
        ("Verification-Based Trust Engine:", " Create an admin control center where moderators manually inspect uploaded host documents (Aadhar card, utility bill, and sales deed) and log every action to an immutable AuditLog collection."),
        ("Double-Booking Collision Avoidance:", " Implement backend date validators to check requested check-in/out ranges against active bookings, automatically blocking overlaps and auto-rejecting conflicting requests when a booking is accepted."),
        ("Exclusive Women's Safety Ratings:", " Build a rating system that restricts safety reviews (covering lighting, locks, host conduct, transport) strictly to female users with completed and validated stays."),
        ("Eco-Friendly Feature Disclosures:", " Enable hosts to check-list active green features (solar energy, waste sorting, rainwater harvesting, energy appliances, public transit, water fixtures) and display them as badges on listing details."),
        ("Community Services Marketplace:", " Help local providers list drivers, tour guides, and house helpers categorized by city/state, making it easy for guests to access nearby amenities.")
    ]
    for title, desc in objectives:
        p = tfo.add_paragraph()
        p.text = f"✔  {title}{desc}"
        p.font.name = "Calibri"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(7)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = TEAL

    # SLIDE 5: Scope of the Project (Detailed)
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide5, LIGHT_GRAY)
    add_header(slide5, "Scope & Operational Boundaries")

    # Column 1: In Scope
    scope_col1 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), Inches(5.7), Inches(5.3))
    scope_col1.fill.solid()
    scope_col1.fill.fore_color.rgb = WHITE
    scope_col1.line.fill.background()

    scope_box1 = slide5.shapes.add_textbox(Inches(0.95), Inches(1.7), Inches(5.3), Inches(4.9))
    tfs1 = scope_box1.text_frame
    tfs1.word_wrap = True

    ps1_h = tfs1.paragraphs[0]
    ps1_h.text = "In-Scope Modules & Roles"
    ps1_h.font.name = "Calibri"
    ps1_h.font.size = Pt(18)
    ps1_h.font.bold = True
    ps1_h.font.color.rgb = DARK_BLUE
    ps1_h.space_after = Pt(12)

    scopes_in = [
        ("Guest Operations:", " Registering accounts, searching listings by city, selecting checkin-checkout dates, submitting bookings, writing reviews, and adding female safety scores (if female)."),
        ("Host Operations:", " Creating normal/shared listing profiles, uploading images and verification PDFs, managing reservation status (accept, reject, cancel), and advertising local marketplace services."),
        ("Admin Portal:", " Reviewing pending properties, checking credentials, publishing verified listings, and checking immutable system logs.")
    ]
    for title, desc in scopes_in:
        p = tfs1.add_paragraph()
        p.text = f"• {title}{desc}"
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = DARK_BLUE

    # Column 2: Out of Scope
    scope_col2 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.883), Inches(1.5), Inches(5.7), Inches(5.3))
    scope_col2.fill.solid()
    scope_col2.fill.fore_color.rgb = WHITE
    scope_col2.line.fill.background()

    scope_box2 = slide5.shapes.add_textbox(Inches(7.083), Inches(1.7), Inches(5.3), Inches(4.9))
    tfs2 = scope_box2.text_frame
    tfs2.word_wrap = True

    ps2_h = tfs2.paragraphs[0]
    ps2_h.text = "System Boundaries & Limitations"
    ps2_h.font.name = "Calibri"
    ps2_h.font.size = Pt(18)
    ps2_h.font.bold = True
    ps2_h.font.color.rgb = DARK_BLUE
    ps2_h.space_after = Pt(12)

    scopes_out = [
        ("Simulated Payment Flow:", " Financial transactions are simulated using pending-accepted status changes; no direct payment gateway (Stripe/Razorpay) is integrated."),
        ("Manual Vetting:", " Document vetting is performed manually by admins through a workspace panel; automated OCR analysis is excluded."),
        ("Local Storage:", " Uploaded documents and listings are saved locally on the hosting server; external cloud servers (S3) are not used."),
        ("No Live Chat:", " Host-guest interactions occur via dashboard booking status triggers; instant chat is excluded.")
    ]
    for title, desc in scopes_out:
        p = tfs2.add_paragraph()
        p.text = f"✖ {title}{desc}"
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = ACCENT_RED

    # SLIDE 6: Module Specification - Part 1 (Detailed Core System)
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide6, LIGHT_GRAY)
    add_header(slide6, "Module Specification (Part 1)")

    # Columns for 3 modules
    mod_widths = Inches(3.7)
    mod_height = Inches(5.3)
    mod_tops = Inches(1.5)
    left_positions = [Inches(0.75), Inches(4.816), Inches(8.883)]

    # Module A: Authentication & Security
    card_a = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_positions[0], mod_tops, mod_widths, mod_height)
    card_a.fill.solid()
    card_a.fill.fore_color.rgb = WHITE
    card_a.line.fill.background()

    box_a = slide6.shapes.add_textbox(left_positions[0]+Inches(0.15), mod_tops+Inches(0.15), mod_widths-Inches(0.3), mod_height-Inches(0.3))
    tfa = box_a.text_frame
    tfa.word_wrap = True
    
    pa_h = tfa.paragraphs[0]
    pa_h.text = "1. Security & Session"
    pa_h.font.name = "Calibri"
    pa_h.font.size = Pt(17)
    pa_h.font.bold = True
    pa_h.font.color.rgb = DARK_BLUE
    pa_h.space_after = Pt(10)

    auth_bullets = [
        "Passport.js Local Strategy is used to manage host/guest credentials.",
        "Uses PBKDF2 password hashing with custom iterations to prevent brute-force attacks.",
        "express-session creates cookie-based sessions, storing login state.",
        "Route Guards: prevents non-logged users from booking, non-owners from editing, and non-admins from vetting.",
        "Local variable middleware dynamically changes navbar links (Login/Signup vs. Logout) based on session."
    ]
    for bullet in auth_bullets:
        p = tfa.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)

    # Module B: Listing CRUD Desk
    card_b = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_positions[1], mod_tops, mod_widths, mod_height)
    card_b.fill.solid()
    card_b.fill.fore_color.rgb = WHITE
    card_b.line.fill.background()

    box_b = slide6.shapes.add_textbox(left_positions[1]+Inches(0.15), mod_tops+Inches(0.15), mod_widths-Inches(0.3), mod_height-Inches(0.3))
    tfb = box_b.text_frame
    tfb.word_wrap = True
    
    pb_h = tfb.paragraphs[0]
    pb_h.text = "2. Listing CRUD Desk"
    pb_h.font.name = "Calibri"
    pb_h.font.size = Pt(17)
    pb_h.font.bold = True
    pb_h.font.color.rgb = DARK_BLUE
    pb_h.space_after = Pt(10)

    listing_bullets = [
        "Properties are listed with titles, categories (normal/shared), prices, meal types, and locations.",
        "Multer middleware handles multipart image and PDF document uploads.",
        "Host uploads government ID (Aadhar) and property proof (utility bill/deed) on creation.",
        "New or edited properties are set to pending and hidden until verified.",
        "Editing a live property resets status to pending, forcing re-vetting to prevent sneaky listings changes."
    ]
    for bullet in listing_bullets:
        p = tfb.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)

    # Module C: Booking Conflict Engine
    card_c = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_positions[2], mod_tops, mod_widths, mod_height)
    card_c.fill.solid()
    card_c.fill.fore_color.rgb = WHITE
    card_c.line.fill.background()

    box_c = slide6.shapes.add_textbox(left_positions[2]+Inches(0.15), mod_tops+Inches(0.15), mod_widths-Inches(0.3), mod_height-Inches(0.3))
    tfc = box_c.text_frame
    tfc.word_wrap = True
    
    pc_h = tfc.paragraphs[0]
    pc_h.text = "3. Booking Engine"
    pc_h.font.name = "Calibri"
    pc_h.font.size = Pt(17)
    pc_h.font.bold = True
    pc_h.font.color.rgb = DARK_BLUE
    pc_h.space_after = Pt(10)

    booking_bullets = [
        "Guests select dates using Flatpickr calendar. Total price, nights, and guests are dynamically calculated.",
        "Collision prevention engine runs checks before reserving rooms.",
        "Validates date overlaps: checks if any accepted/pending booking exists where checkIn < requestedCheckOut and checkOut > requestedCheckIn.",
        "Accepting a booking locks the dates, auto-rejecting conflicting pending requests.",
        "Rejections or cancellations release the dates back to the available pool."
    ]
    for bullet in booking_bullets:
        p = tfc.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)

    # SLIDE 7: Module Specification - Part 2 (Detailed Features)
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide7, LIGHT_GRAY)
    add_header(slide7, "Module Specification (Part 2)")

    # Columns for 4 modules
    mod_widths4 = Inches(2.7)
    left_positions4 = [Inches(0.75), Inches(3.7), Inches(6.65), Inches(9.6)]

    # Mod 1: Women's Safety
    card_m1 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_positions4[0], mod_tops, mod_widths4, mod_height)
    card_m1.fill.solid()
    card_m1.fill.fore_color.rgb = WHITE
    card_m1.line.fill.background()

    box_m1 = slide7.shapes.add_textbox(left_positions4[0]+Inches(0.1), mod_tops+Inches(0.15), mod_widths4-Inches(0.2), mod_height-Inches(0.3))
    tf_m1 = box_m1.text_frame
    tf_m1.word_wrap = True
    pm1_h = tf_m1.paragraphs[0]
    pm1_h.text = "4. Women Safety"
    pm1_h.font.name = "Calibri"
    pm1_h.font.size = Pt(16)
    pm1_h.font.bold = True
    pm1_h.font.color.rgb = DARK_BLUE
    pm1_h.space_after = Pt(8)

    ws_b = [
        "Restricted access: Checks guest.gender === 'female' in database.",
        "Completed stay check: Only allows reviews for accepted checkout dates in the past.",
        "Collects binary tags: lighting, locks/privacy, host conduct, transport, neighborhood.",
        "Uses MongoDB aggregate pipelines to calculate averages.",
        "Displays safety ratings separately to guide female travelers."
    ]
    for bullet in ws_b:
        p = tf_m1.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = "Calibri"
        p.font.size = Pt(11.5)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(5)

    # Mod 2: Eco-Friendly Disclosures (Modified to remove scores)
    card_m2 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_positions4[1], mod_tops, mod_widths4, mod_height)
    card_m2.fill.solid()
    card_m2.fill.fore_color.rgb = WHITE
    card_m2.line.fill.background()

    box_m2 = slide7.shapes.add_textbox(left_positions4[1]+Inches(0.1), mod_tops+Inches(0.15), mod_widths4-Inches(0.2), mod_height-Inches(0.3))
    tf_m2 = box_m2.text_frame
    tf_m2.word_wrap = True
    pm2_h = tf_m2.paragraphs[0]
    pm2_h.text = "5. Eco Disclosures"
    pm2_h.font.name = "Calibri"
    pm2_h.font.size = Pt(16)
    pm2_h.font.bold = True
    pm2_h.font.color.rgb = DARK_BLUE
    pm2_h.space_after = Pt(8)

    sust_b = [
        "Hosts can check-list which eco features are active at their property.",
        "Checks cover solar energy, recycling, rainwater harvesting, efficient devices, transport, water saving.",
        "Lists active features as informative badges directly on property details.",
        "NO mathematical rating or calculated scores are applied.",
        "Empowers travelers to evaluate green amenities transparently."
    ]
    for bullet in sust_b:
        p = tf_m2.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = "Calibri"
        p.font.size = Pt(11.5)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(5)

    # Mod 3: Local Services
    card_m3 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_positions4[2], mod_tops, mod_widths4, mod_height)
    card_m3.fill.solid()
    card_m3.fill.fore_color.rgb = WHITE
    card_m3.line.fill.background()

    box_m3 = slide7.shapes.add_textbox(left_positions4[2]+Inches(0.1), mod_tops+Inches(0.15), mod_widths4-Inches(0.2), mod_height-Inches(0.3))
    tf_m3 = box_m3.text_frame
    tf_m3.word_wrap = True
    pm3_h = tf_m3.paragraphs[0]
    pm3_h.text = "6. Local Services"
    pm3_h.font.name = "Calibri"
    pm3_h.font.size = Pt(16)
    pm3_h.font.bold = True
    pm3_h.font.color.rgb = DARK_BLUE
    pm3_h.space_after = Pt(8)

    serv_b = [
        "Directory of drivers, guides, and house helpers.",
        "Help local providers list services to generate supplemental income.",
        "Searchable database filtered by city parameters.",
        "Unique phone check prevents duplicates and spam.",
        "Integrates with listings, linking local support directly."
    ]
    for bullet in serv_b:
        p = tf_m3.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = "Calibri"
        p.font.size = Pt(11.5)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(5)

    # Mod 4: Admin Vetting & Audits
    card_m4 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_positions4[3], mod_tops, mod_widths4, mod_height)
    card_m4.fill.solid()
    card_m4.fill.fore_color.rgb = WHITE
    card_m4.line.fill.background()

    box_m4 = slide7.shapes.add_textbox(left_positions4[3]+Inches(0.1), mod_tops+Inches(0.15), mod_widths4-Inches(0.2), mod_height-Inches(0.3))
    tf_m4 = box_m4.text_frame
    tf_m4.word_wrap = True
    pm4_h = tf_m4.paragraphs[0]
    pm4_h.text = "7. Vetting & Audits"
    pm4_h.font.name = "Calibri"
    pm4_h.font.size = Pt(16)
    pm4_h.font.bold = True
    pm4_h.font.color.rgb = DARK_BLUE
    pm4_h.space_after = Pt(8)

    audit_b = [
        "Vetting queue displays unverified properties.",
        "Vetting workspace to review uploaded documents.",
        "Approves listings or enters rejection reasons.",
        "Writes immutable event logs to AuditLog collection.",
        "Audit Log tracks: admin id, listing id, action, rejection reasons, and timestamps."
    ]
    for bullet in audit_b:
        p = tf_m4.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = "Calibri"
        p.font.size = Pt(11.5)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(5)

    # SLIDE 8: Project Diagrams - System Architecture
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide8, LIGHT_GRAY)
    add_header(slide8, "MVC System Architecture")

    # Column 1: Description Text
    col8_card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), Inches(4.5), Inches(5.3))
    col8_card.fill.solid()
    col8_card.fill.fore_color.rgb = WHITE
    col8_card.line.fill.background()

    desc_box8 = slide8.shapes.add_textbox(Inches(0.95), Inches(1.7), Inches(4.1), Inches(4.9))
    tf8 = desc_box8.text_frame
    tf8.word_wrap = True
    p8_h = tf8.paragraphs[0]
    p8_h.text = "MVC Pattern Execution"
    p8_h.font.name = "Calibri"
    p8_h.font.size = Pt(20)
    p8_h.font.bold = True
    p8_h.font.color.rgb = DARK_BLUE
    p8_h.space_after = Pt(12)

    desc8_pts = [
        "Browser (Frontend Client): Sends HTTP requests (POST, GET, PUT) for logins, listings, and bookings. Interacts with EJS templates containing Flatpickr calendars.",
        "Middleware Routing (Guards): Intercepts requests, checking session cookies via Passport.js to verify active logins, listing ownership, and admin access.",
        "Controllers (Backend Engine): Executes business calculations, processes uploads, runs overlap validators, and fetches collections.",
        "Model / DB (Mongoose & MongoDB): Defines schemas and indexes for Users, Listings, Bookings, Services, and Logs."
    ]
    for pt in desc8_pts:
        p = tf8.add_paragraph()
        p.text = f"• {pt}"
        p.font.name = "Calibri"
        p.font.size = Pt(12.5)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)

    # Column 2: Large Placeholder Box for Architecture Diagram
    diag8_card = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.6), Inches(1.5), Inches(7.0), Inches(5.3))
    diag8_card.fill.solid()
    diag8_card.fill.fore_color.rgb = LIGHT_BLUE
    diag8_card.line.color.rgb = TEAL
    diag8_card.line.width = Pt(1.5)

    diag8_box = slide8.shapes.add_textbox(Inches(5.7), Inches(3.0), Inches(6.8), Inches(2.0))
    tf8_d = diag8_box.text_frame
    tf8_d.word_wrap = True
    p8_d = tf8_d.paragraphs[0]
    p8_d.alignment = PP_ALIGN.CENTER
    p8_d.text = "[ PLACEHOLDER: SYSTEM ARCHITECTURE DIAGRAM ]"
    p8_d.font.name = "Calibri"
    p8_d.font.size = Pt(20)
    p8_d.font.bold = True
    p8_d.font.color.rgb = DARK_BLUE
    p8_d.space_after = Pt(10)

    p8_d2 = tf8_d.add_paragraph()
    p8_d2.alignment = PP_ALIGN.CENTER
    p8_d2.text = "(Copy and paste the MVC Architecture flowchart from Chapter 2 of the Project Report here)"
    p8_d2.font.name = "Calibri"
    p8_d2.font.size = Pt(13)
    p8_d2.font.color.rgb = DARK_GRAY

    # SLIDE 9: Project Diagrams - Activity Flow
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide9, LIGHT_GRAY)
    add_header(slide9, "System Activity Diagram Flow")

    # Column 1: Description Text
    col9_card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), Inches(4.5), Inches(5.3))
    col9_card.fill.solid()
    col9_card.fill.fore_color.rgb = WHITE
    col9_card.line.fill.background()

    desc_box9 = slide9.shapes.add_textbox(Inches(0.95), Inches(1.7), Inches(4.1), Inches(4.9))
    tf9 = desc_box9.text_frame
    tf9.word_wrap = True
    p9_h = tf9.paragraphs[0]
    p9_h.text = "Activity Steps Explained"
    p9_h.font.name = "Calibri"
    p9_h.font.size = Pt(20)
    p9_h.font.bold = True
    p9_h.font.color.rgb = DARK_BLUE
    p9_h.space_after = Pt(12)

    desc9_pts = [
        "1. Login/Signup check: Validates credentials, directs users to Guest, Host, or Admin dashboard.",
        "2. Guest pipeline: Searches listings -> selects dates -> check conflicts -> creates pending booking request -> checks out -> female rating (if female).",
        "3. Host pipeline: Fills details -> checks eco-features -> uploads docs -> waits for vetting -> live status -> approves or rejects guest booking request.",
        "4. Admin pipeline: Vets docs -> approves/rejects listing -> logs results to AuditLog collection."
    ]
    for pt in desc9_pts:
        p = tf9.add_paragraph()
        p.text = pt
        p.font.name = "Calibri"
        p.font.size = Pt(12.5)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)

    # Column 2: Large Placeholder Box for Activity Diagram
    diag9_card = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.6), Inches(1.5), Inches(7.0), Inches(5.3))
    diag9_card.fill.solid()
    diag9_card.fill.fore_color.rgb = LIGHT_BLUE
    diag9_card.line.color.rgb = TEAL
    diag9_card.line.width = Pt(1.5)

    diag9_box = slide9.shapes.add_textbox(Inches(5.7), Inches(3.0), Inches(6.8), Inches(2.0))
    tf9_d = diag9_box.text_frame
    tf9_d.word_wrap = True
    p9_d = tf9_d.paragraphs[0]
    p9_d.alignment = PP_ALIGN.CENTER
    p9_d.text = "[ PLACEHOLDER: SYSTEM ACTIVITY FLOWCHART ]"
    p9_d.font.name = "Calibri"
    p9_d.font.size = Pt(20)
    p9_d.font.bold = True
    p9_d.font.color.rgb = DARK_BLUE
    p9_d.space_after = Pt(10)

    p9_d2 = tf9_d.add_paragraph()
    p9_d2.alignment = PP_ALIGN.CENTER
    p9_d2.text = "(Copy and paste the detailed multi-swimlane Activity flowchart from Chapter 2 of the Project Report)"
    p9_d2.font.name = "Calibri"
    p9_d2.font.size = Pt(13)
    p9_d2.font.color.rgb = DARK_GRAY

    # SLIDE 10: Project Diagrams - Component Sequence Flows (Detailed)
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide10, LIGHT_GRAY)
    add_header(slide10, "Component Sequence Flows")

    # Column 1: Description Text
    col10_card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), Inches(4.5), Inches(5.3))
    col10_card.fill.solid()
    col10_card.fill.fore_color.rgb = WHITE
    col10_card.line.fill.background()

    desc_box10 = slide10.shapes.add_textbox(Inches(0.95), Inches(1.7), Inches(4.1), Inches(4.9))
    tf10 = desc_box10.text_frame
    tf10.word_wrap = True
    p10_h = tf10.paragraphs[0]
    p10_h.text = "Sequence logic step-by-step"
    p10_h.font.name = "Calibri"
    p10_h.font.size = Pt(20)
    p10_h.font.bold = True
    p10_h.font.color.rgb = DARK_BLUE
    p10_h.space_after = Pt(12)

    desc10_pts = [
        "A. Vetting flow: Host submits listing + uploads Aadhar/utility bill. Multer processes and saves files locally under /public/uploads/. The listing is created with status 'pending' in MongoDB. The Admin views the listing in the queue, approves it, and an immutable log entry is generated in the database.",
        "B. Booking flow: Guest requests booking. Booking controller runs conflict checks. If dates are free, booking is saved as pending. If the Host accepts, dates are locked, and overlapping bookings are auto-rejected."
    ]
    for pt in desc10_pts:
        p = tf10.add_paragraph()
        p.text = pt
        p.font.name = "Calibri"
        p.font.size = Pt(12.5)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)

    # Column 2: Large Placeholder Box for Sequence Diagrams
    diag10_card = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.6), Inches(1.5), Inches(7.0), Inches(5.3))
    diag10_card.fill.solid()
    diag10_card.fill.fore_color.rgb = LIGHT_BLUE
    diag10_card.line.color.rgb = TEAL
    diag10_card.line.width = Pt(1.5)

    diag10_box = slide10.shapes.add_textbox(Inches(5.7), Inches(3.0), Inches(6.8), Inches(2.0))
    tf10_d = diag10_box.text_frame
    tf10_d.word_wrap = True
    p10_d = tf10_d.paragraphs[0]
    p10_d.alignment = PP_ALIGN.CENTER
    p10_d.text = "[ PLACEHOLDER: SYSTEM SEQUENCE FLOWS ]"
    p10_d.font.name = "Calibri"
    p10_d.font.size = Pt(20)
    p10_d.font.bold = True
    p10_d.font.color.rgb = DARK_BLUE
    p10_d.space_after = Pt(10)

    p10_d2 = tf10_d.add_paragraph()
    p10_d2.alignment = PP_ALIGN.CENTER
    p10_d2.text = "(Copy and paste the Vetting Sequence diagram and Booking Conflict sequence diagram from Chapter 2 of the Project Report)"
    p10_d2.font.name = "Calibri"
    p10_d2.font.size = Pt(13)
    p10_d2.font.color.rgb = DARK_GRAY

    # SLIDE 11: User Interface Screens - Part 1
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide11, LIGHT_GRAY)
    add_header(slide11, "User Interface: Discovery & Showcases")

    # Column 1 Placeholder
    c11_card1 = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.5), Inches(5.7), Inches(5.3))
    c11_card1.fill.solid()
    c11_card1.fill.fore_color.rgb = LIGHT_BLUE
    c11_card1.line.color.rgb = TEAL
    c11_card1.line.width = Pt(1.5)

    box11_1 = slide11.shapes.add_textbox(Inches(0.85), Inches(2.8), Inches(5.5), Inches(2.5))
    tf11_1 = box11_1.text_frame
    tf11_1.word_wrap = True
    p11_1 = tf11_1.paragraphs[0]
    p11_1.alignment = PP_ALIGN.CENTER
    p11_1.text = "[ PLACEHOLDER: HOMEPAGE & DISCOVERY ]"
    p11_1.font.name = "Calibri"
    p11_1.font.size = Pt(18)
    p11_1.font.bold = True
    p11_1.font.color.rgb = DARK_BLUE
    p11_1.space_after = Pt(10)

    p11_1b = tf11_1.add_paragraph()
    p11_1b.alignment = PP_ALIGN.CENTER
    p11_1b.text = "Insert screenshot of index.ejs showing:\n- Navigation Header & Search Filter Bar\n- Shared Category Filters\n- Property Listing Cards with Eco-Feature Badges"
    p11_1b.font.name = "Calibri"
    p11_1b.font.size = Pt(12)
    p11_1b.font.color.rgb = DARK_GRAY

    # Column 2 Placeholder
    c11_card2 = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.883), Inches(1.5), Inches(5.7), Inches(5.3))
    c11_card2.fill.solid()
    c11_card2.fill.fore_color.rgb = LIGHT_BLUE
    c11_card2.line.color.rgb = TEAL
    c11_card2.line.width = Pt(1.5)

    box11_2 = slide11.shapes.add_textbox(Inches(6.983), Inches(2.8), Inches(5.5), Inches(2.5))
    tf11_2 = box11_2.text_frame
    tf11_2.word_wrap = True
    p11_2 = tf11_2.paragraphs[0]
    p11_2.alignment = PP_ALIGN.CENTER
    p11_2.text = "[ PLACEHOLDER: PROPERTY DETAIL VIEW ]"
    p11_2.font.name = "Calibri"
    p11_2.font.size = Pt(18)
    p11_2.font.bold = True
    p11_2.font.color.rgb = DARK_BLUE
    p11_2.space_after = Pt(10)

    p11_2b = tf11_2.add_paragraph()
    p11_2b.alignment = PP_ALIGN.CENTER
    p11_2b.text = "Insert screenshot of show.ejs showing:\n- Image Carousel & Description text\n- Flatpickr check-in/out form\n- Sustainability checklists & Women's safety ratings"
    p11_2b.font.name = "Calibri"
    p11_2b.font.size = Pt(12)
    p11_2b.font.color.rgb = DARK_GRAY

    # SLIDE 12: User Interface Screens - Part 2
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide12, LIGHT_GRAY)
    add_header(slide12, "User Interface: Portals & Marketplace")

    # Column 1 Placeholder
    c12_card1 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.5), Inches(5.7), Inches(5.3))
    c12_card1.fill.solid()
    c12_card1.fill.fore_color.rgb = LIGHT_BLUE
    c12_card1.line.color.rgb = TEAL
    c12_card1.line.width = Pt(1.5)

    box12_1 = slide12.shapes.add_textbox(Inches(0.85), Inches(2.8), Inches(5.5), Inches(2.5))
    tf12_1 = box12_1.text_frame
    tf12_1.word_wrap = True
    p12_1 = tf12_1.paragraphs[0]
    p12_1.alignment = PP_ALIGN.CENTER
    p12_1.text = "[ PLACEHOLDER: BOOKINGS & STATE CONTROL ]"
    p12_1.font.name = "Calibri"
    p12_1.font.size = Pt(18)
    p12_1.font.bold = True
    p12_1.font.color.rgb = DARK_BLUE
    p12_1.space_after = Pt(10)

    p12_1b = tf12_1.add_paragraph()
    p12_1b.alignment = PP_ALIGN.CENTER
    p12_1b.text = "Insert screenshot of my-bookings.ejs showing:\n- Guest Bookings checklist and receipt download\n- Host Reservations list dashboard\n- Booking state buttons (Accept, Reject, Cancel)"
    p12_1b.font.name = "Calibri"
    p12_1b.font.size = Pt(12)
    p12_1b.font.color.rgb = DARK_GRAY

    # Column 2 Placeholder
    c12_card2 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.883), Inches(1.5), Inches(5.7), Inches(5.3))
    c12_card2.fill.solid()
    c12_card2.fill.fore_color.rgb = LIGHT_BLUE
    c12_card2.line.color.rgb = TEAL
    c12_card2.line.width = Pt(1.5)

    box12_2 = slide12.shapes.add_textbox(Inches(6.983), Inches(2.8), Inches(5.5), Inches(2.5))
    tf12_2 = box12_2.text_frame
    tf12_2.word_wrap = True
    p12_2 = tf12_2.paragraphs[0]
    p12_2.alignment = PP_ALIGN.CENTER
    p12_2.text = "[ PLACEHOLDER: SERVICES MARKETPLACE ]"
    p12_2.font.name = "Calibri"
    p12_2.font.size = Pt(18)
    p12_2.font.bold = True
    p12_2.font.color.rgb = DARK_BLUE
    p12_2.space_after = Pt(10)

    p12_2b = tf12_2.add_paragraph()
    p12_2b.alignment = PP_ALIGN.CENTER
    p12_2b.text = "Insert screenshot of services/index.ejs showing:\n- Category listings (Drivers, Guides, Helpers)\n- Local contact directory filtered by city\n- Registration form for new local service providers"
    p12_2b.font.name = "Calibri"
    p12_2b.font.size = Pt(12)
    p12_2b.font.color.rgb = DARK_GRAY

    # SLIDE 13: User Interface Screens - Part 3
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide13, LIGHT_GRAY)
    add_header(slide13, "User Interface: Admin Console")

    # Column 1 Placeholder
    c13_card1 = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.5), Inches(5.7), Inches(5.3))
    c13_card1.fill.solid()
    c13_card1.fill.fore_color.rgb = LIGHT_BLUE
    c13_card1.line.color.rgb = TEAL
    c13_card1.line.width = Pt(1.5)

    box13_1 = slide13.shapes.add_textbox(Inches(0.85), Inches(2.8), Inches(5.5), Inches(2.5))
    tf13_1 = box13_1.text_frame
    tf13_1.word_wrap = True
    p13_1 = tf13_1.paragraphs[0]
    p13_1.alignment = PP_ALIGN.CENTER
    p13_1.text = "[ PLACEHOLDER: ADMIN VETTING QUEUE ]"
    p13_1.font.name = "Calibri"
    p13_1.font.size = Pt(18)
    p13_1.font.bold = True
    p13_1.font.color.rgb = DARK_BLUE
    p13_1.space_after = Pt(10)

    p13_1b = tf13_1.add_paragraph()
    p13_1b.alignment = PP_ALIGN.CENTER
    p13_1b.text = "Insert screenshot of admin/dashboard.ejs showing:\n- Pending host properties vetting queue\n- Government document inspect dashboard\n- Approve / Reject buttons with rejection inputs"
    p13_1b.font.name = "Calibri"
    p13_1b.font.size = Pt(12)
    p13_1b.font.color.rgb = DARK_GRAY

    # Column 2 Placeholder
    c13_card2 = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.883), Inches(1.5), Inches(5.7), Inches(5.3))
    c13_card2.fill.solid()
    c13_card2.fill.fore_color.rgb = LIGHT_BLUE
    c13_card2.line.color.rgb = TEAL
    c13_card2.line.width = Pt(1.5)

    box13_2 = slide13.shapes.add_textbox(Inches(6.983), Inches(2.8), Inches(5.5), Inches(2.5))
    tf13_2 = box13_2.text_frame
    tf13_2.word_wrap = True
    p13_2 = tf13_2.paragraphs[0]
    p13_2.alignment = PP_ALIGN.CENTER
    p13_2.text = "[ PLACEHOLDER: SYSTEM AUDIT LOGS ]"
    p13_2.font.name = "Calibri"
    p13_2.font.size = Pt(18)
    p13_2.font.bold = True
    p13_2.font.color.rgb = DARK_BLUE
    p13_2.space_after = Pt(10)

    p13_2b = tf13_2.add_paragraph()
    p13_2b.alignment = PP_ALIGN.CENTER
    p13_2b.text = "Insert screenshot of audit trails table showing:\n- Active Admin identifier\n- Affected Listing details\n- Performed Actions (Approve/Reject), Rejection Reasons, and Timestamps"
    p13_2b.font.name = "Calibri"
    p13_2b.font.size = Pt(12)
    p13_2b.font.color.rgb = DARK_GRAY

    # SLIDE 14: Test Plan (Detailed Scenario Explanations)
    slide14 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide14, LIGHT_GRAY)
    add_header(slide14, "Testing Plan & Quality Assurance")

    test_card = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.3))
    test_card.fill.solid()
    test_card.fill.fore_color.rgb = WHITE
    test_card.line.fill.background()

    test_box = slide14.shapes.add_textbox(Inches(0.95), Inches(1.6), Inches(11.333), Inches(5.1))
    tft = test_box.text_frame
    tft.word_wrap = True

    pt_h = tft.paragraphs[0]
    pt_h.text = "Manual Testing Scenarios & Verification Methods"
    pt_h.font.name = "Calibri"
    pt_h.font.size = Pt(18)
    pt_h.font.bold = True
    pt_h.font.color.rgb = DARK_BLUE
    pt_h.space_after = Pt(10)

    test_scenarios = [
        ("Authentication Verification (TC-AUTH):", " Checks form validators (non-empty username, valid email formats), enforces password security checks (Passport.js), and verifies session creation and login redirections (9 scenarios)."),
        ("Listings Integrity (TC-LST):", " Checks that listing forms block invalid entries (negative prices, empty descriptions), verifies owner-only edit controls, and confirms search filtering options by city names (7 scenarios)."),
        ("Admin Verification Queue (TC-ADM):", " Checks that host Aadhar cards, deeds, and bills are visible in the vetting dashboard. Confirms approving/rejecting listings logs actions to MongoDB AuditLog (5 scenarios)."),
        ("Guest Booking Pipeline (TC-BKG):", " Verifies checkin-checkout calendar ranges. Confirms the collision checker blocks bookings that overlap with existing accepted bookings (5 scenarios)."),
        ("Host Booking State (TC-HMG):", " Verifies reservation management. Confirms accepting a booking locks the selected dates and auto-rejects overlapping requests (5 scenarios)."),
        ("Local Services Marketplace (TC-SRV):", " Enforces unique contact phone constraints to prevent spam and checks that directory searches correctly filter by city parameters (6 scenarios).")
    ]
    for title, desc in test_scenarios:
        p = tft.add_paragraph()
        p.text = f"✔  {title}{desc}"
        p.font.name = "Calibri"
        p.font.size = Pt(12.5)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = TEAL

    # SLIDE 15: Test Plan - Testing Results Metrics (Detailed Bug Fixes)
    slide15 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide15, LIGHT_GRAY)
    add_header(slide15, "Testing Metrics & Bug Resolutions")

    # Column 1: Metrics
    col15_card1 = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), Inches(5.7), Inches(5.3))
    col15_card1.fill.solid()
    col15_card1.fill.fore_color.rgb = WHITE
    col15_card1.line.fill.background()

    metrics_box = slide15.shapes.add_textbox(Inches(0.95), Inches(1.7), Inches(5.3), Inches(4.9))
    tf15_1 = metrics_box.text_frame
    tf15_1.word_wrap = True

    p15_1h = tf15_1.paragraphs[0]
    p15_1h.text = "Manual Test Metrics Summary"
    p15_1h.font.name = "Calibri"
    p15_1h.font.size = Pt(18)
    p15_1h.font.bold = True
    p15_1h.font.color.rgb = DARK_BLUE
    p15_1h.space_after = Pt(15)

    metrics = [
        "1. Authentication: 9 Runs, 9 Passed, 0 Failed (100% Pass)",
        "2. Listings CRUD: 7 Runs, 7 Passed, 0 Failed (100% Pass)",
        "3. Admin Vetting: 5 Runs, 5 Passed, 0 Failed (100% Pass)",
        "4. Guest Bookings: 5 Runs, 5 Passed, 0 Failed (100% Pass)",
        "5. Host Management: 5 Runs, 5 Passed, 0 Failed (100% Pass)",
        "6. Services Directory: 6 Runs, 6 Passed, 0 Failed (100% Pass)",
        "Total Execution: 37 Cases, 37 Passed, 0 Failed (100% Pass)"
    ]
    for metric in metrics:
        p = tf15_1.add_paragraph()
        p.text = metric
        p.font.name = "Calibri"
        p.font.size = Pt(12.5)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
        if "Total" in metric:
            p.font.bold = True
            p.font.color.rgb = TEAL

    # Column 2: Defect Resolution
    col15_card2 = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.883), Inches(1.5), Inches(5.7), Inches(5.3))
    col15_card2.fill.solid()
    col15_card2.fill.fore_color.rgb = WHITE
    col15_card2.line.fill.background()

    defects_box = slide15.shapes.add_textbox(Inches(7.083), Inches(1.7), Inches(5.3), Inches(4.9))
    tf15_2 = defects_box.text_frame
    tf15_2.word_wrap = True

    p15_2h = tf15_2.paragraphs[0]
    p15_2h.text = "Fixed Defects & Solutions"
    p15_2h.font.name = "Calibri"
    p15_2h.font.size = Pt(18)
    p15_2h.font.bold = True
    p15_2h.font.color.rgb = DARK_BLUE
    p15_2h.space_after = Pt(15)

    defects = [
        ("Listing Edit Location Bug:", " Editing a property did not save city, state, or country fields. This occurred because HTML input tags in edit.ejs were missing the 'name' attributes. Added name attributes to restore data binding."),
        ("Duplicate Service Contacts:", " Service providers could list duplicate phone numbers, causing spam. Added unique schema indexes and modified the controller to catch error code 11000 and return custom flash alerts."),
        ("Date Format Selection Bug:", " Calendar inputs returned inconsistent date strings across different browsers, breaking backend date queries. Standardized input formats to YYYY-MM-DD using Flatpickr settings.")
    ]
    for title, desc in defects:
        p = tf15_2.add_paragraph()
        p.text = f"• {title}{desc}"
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = DARK_BLUE

    # SLIDE 16: Conclusion & Future Scope (Detailed Explanation)
    slide16 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide16, LIGHT_GRAY)
    add_header(slide16, "Conclusion & Future Enhancements")

    # Column 1: Summary Conclusion
    col16_card1 = slide16.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), Inches(5.7), Inches(5.3))
    col16_card1.fill.solid()
    col16_card1.fill.fore_color.rgb = WHITE
    col16_card1.line.fill.background()

    concl_box = slide16.shapes.add_textbox(Inches(0.95), Inches(1.7), Inches(5.3), Inches(4.9))
    tf16_1 = concl_box.text_frame
    tf16_1.word_wrap = True

    p16_1h = tf16_1.paragraphs[0]
    p16_1h.text = "Project Conclusion"
    p16_1h.font.name = "Calibri"
    p16_1h.font.size = Pt(18)
    p16_1h.font.bold = True
    p16_1h.font.color.rgb = DARK_BLUE
    p16_1h.space_after = Pt(15)

    concl_text = (
        "Nestify is a community-focused web application built on the MERN stack. "
        "It successfully implements a secure shared-accommodation booking system "
        "with several innovative community modules:\n\n"
        "1. Credential Vetting: Reduces fraud by requiring host ID and property document uploads before listing.\n"
        "2. Verified Safety Scores: Restricts safety reviews to female guests with confirmed stays, providing solo female travelers with trusted feedback.\n"
        "3. Eco-Friendly Disclosures: Displays host-declared eco-features (like solar panels or rainwater harvesting) directly on the details view, aiding transparency.\n"
        "4. Services Directory: Connects local workers with travelers, helping providers earn supplementary income."
    )
    p = tf16_1.add_paragraph()
    p.text = concl_text
    p.font.name = "Calibri"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(10)

    # Column 2: Future Scope
    col16_card2 = slide16.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.883), Inches(1.5), Inches(5.7), Inches(5.3))
    col16_card2.fill.solid()
    col16_card2.fill.fore_color.rgb = WHITE
    col16_card2.line.fill.background()

    future_box = slide16.shapes.add_textbox(Inches(7.083), Inches(1.7), Inches(5.3), Inches(4.9))
    tf16_2 = future_box.text_frame
    tf16_2.word_wrap = True

    p16_2h = tf16_2.paragraphs[0]
    p16_2h.text = "Future Roadmap Plans"
    p16_2h.font.name = "Calibri"
    p16_2h.font.size = Pt(18)
    p16_2h.font.bold = True
    p16_2h.font.color.rgb = DARK_BLUE
    p16_2h.space_after = Pt(15)

    futures = [
        ("Live Payment Processing:", " Integrate Stripe or Razorpay APIs to handle secure escrow booking payments, automated host payouts, and platform commissions."),
        ("Real-time Chat Integrations (Socket.io):", " Replace session state polling with instant messaging between hosts and guests to coordinate arrival/departure details."),
        ("Interactive Geocoding Map (Mapbox):", " Embed maps to show property pins, support search-by-radius, and display nearby service providers on visual layouts."),
        ("Automated Document OCR Analysis:", " Integrate text-recognition APIs to automatically parse and vet uploaded deeds and ID cards, reducing manual admin review overhead."),
        ("Mobile PWA Support:", " Turn the web app into a Progressive Web Application to support offline calendar lookups and app-like mobile styling.")
    ]
    for title, desc in futures:
        p = tf16_2.add_paragraph()
        p.text = f"• {title}{desc}"
        p.font.name = "Calibri"
        p.font.size = Pt(11.5)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = TEAL

    # Save presentation
    output_filename = "Nestify_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation successfully created: {output_filename}")

if __name__ == "__main__":
    build_presentation()
